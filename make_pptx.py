"""
PETase-ML Presentation — clean, large-text, easy to read from a distance.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────
NAVY  = RGBColor(0x0b, 0x16, 0x29)
NAVY2 = RGBColor(0x11, 0x1f, 0x38)
DARK  = RGBColor(0x0d, 0x25, 0x48)
MID   = RGBColor(0x1e, 0x3a, 0x5f)
TEAL  = RGBColor(0x10, 0xb9, 0x81)
BLUE  = RGBColor(0x3b, 0x82, 0xf6)
WHITE = RGBColor(0xf1, 0xf5, 0xf9)
LIGHT = RGBColor(0xcb, 0xd5, 0xe1)
DIM   = RGBColor(0x94, 0xa3, 0xb8)
RED   = RGBColor(0xf8, 0x71, 0x71)
AMBER = RGBColor(0xfb, 0xbf, 0x24)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitive helpers ──────────────────────────────────────
def new_slide(bg=NAVY):
    sl = prs.slides.add_slide(blank)
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = bg
    return sl

def box(sl, text, l, t, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def rect(sl, l, t, w, h, fill=NAVY2, border=MID, border_pt=1.0):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

# ── Compound helpers ───────────────────────────────────────
LABEL_TOP = Inches(0.38)
TITLE_TOP = Inches(0.78)
BODY_TOP  = Inches(2.05)
L = Inches(0.9)   # left margin
R = Inches(11.5)  # right edge of full-width content

def slide_label(sl, text):
    box(sl, text.upper(), L, LABEL_TOP, Inches(11), Inches(0.35),
        size=10, bold=True, color=TEAL)

def slide_title(sl, text, size=42):
    box(sl, text, L, TITLE_TOP, Inches(11.5), Inches(1.25),
        size=size, bold=True, color=WHITE)

def divider_line(sl, top=Inches(2.0)):
    s = sl.shapes.add_shape(1, L, top, Inches(0.55), Pt(5))
    s.fill.solid(); s.fill.fore_color.rgb = TEAL
    s.line.fill.background()

def body_text(sl, text, l=L, t=BODY_TOP, w=Inches(11.5), size=20, color=LIGHT, italic=False):
    box(sl, text, l, t, w, H - t - Inches(0.3), size=size, color=color, italic=italic)

def bullets(sl, items, l=L, t=BODY_TOP, w=Inches(11.5), size=20):
    tb = sl.shapes.add_textbox(l, t, w, H - t - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = "\u2022  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = LIGHT

def big_stat(sl, number, label, l, t, w=Inches(3.5), h=Inches(1.4),
             num_color=TEAL, border=TEAL):
    rect(sl, l, t, w, h, fill=DARK, border=border, border_pt=1.5)
    box(sl, number, l + Inches(0.15), t + Inches(0.1), w - Inches(0.3), Inches(0.72),
        size=36, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    box(sl, label, l + Inches(0.15), t + Inches(0.82), w - Inches(0.3), Inches(0.55),
        size=12, color=DIM, align=PP_ALIGN.CENTER)

def card(sl, heading, body, l, t, w, h, border=TEAL, body_size=17):
    rect(sl, l, t, w, h, fill=NAVY2, border=border, border_pt=2)
    box(sl, heading.upper(), l + Inches(0.22), t + Inches(0.14),
        w - Inches(0.44), Inches(0.32), size=10, bold=True, color=border)
    box(sl, body, l + Inches(0.22), t + Inches(0.5),
        w - Inches(0.44), h - Inches(0.6), size=body_size, color=LIGHT)

def two_col_split(sl):
    """Returns (left_x, right_x, col_width) for a symmetric two-col layout."""
    cw = Inches(5.7)
    return L, Inches(7.0), cw

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════
sl = new_slide()

# Right accent panel
rect(sl, Inches(8.2), Inches(0), Inches(5.13), H,
     fill=RGBColor(0x0d, 0x22, 0x44), border=None)

box(sl, "PETase-ML", L, Inches(1.5), Inches(7.5), Inches(2.0),
    size=80, bold=True, color=WHITE)
box(sl, "Engineering plastic-degrading enzymes\nwith machine learning.",
    L, Inches(3.55), Inches(7.2), Inches(1.0), size=24, color=DIM)

# Pills row
for i, pill in enumerate(["Machine Learning", "Protein Engineering", "Bioremediation"]):
    px = L + i * Inches(3.1)
    rect(sl, px, Inches(4.8), Inches(2.85), Inches(0.45),
         fill=DARK, border=TEAL, border_pt=1)
    box(sl, pill, px + Inches(0.1), Inches(4.83), Inches(2.65), Inches(0.38),
        size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

box(sl, "Los Altos Hacks  \u00b7  2026",
    L, Inches(5.55), Inches(5), Inches(0.4), size=14, color=DIM)

# Right-panel labels
box(sl, "Sequence In.", Inches(8.6), Inches(2.0), Inches(4.5), Inches(0.6),
    size=22, bold=True, color=TEAL)
box(sl, "Ranked Candidates Out.", Inches(8.6), Inches(2.65), Inches(4.5), Inches(0.6),
    size=22, bold=True, color=WHITE)
box(sl, "No crystal structure needed.\nNo software to install.",
    Inches(8.6), Inches(3.45), Inches(4.5), Inches(0.9), size=16, color=DIM)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "The Problem")
slide_title(sl, "380 million tons of plastic per year.\nLess than 10% is ever recycled.")
divider_line(sl)

# Three big stats across the bottom half
y = Inches(2.5)
big_stat(sl, "380M", "metric tons of plastic produced annually",
         L, y, num_color=WHITE, border=BLUE)
big_stat(sl, "~9%", "of plastic ever actually recycled",
         Inches(5.3), y, num_color=RED, border=RED)
big_stat(sl, "200 yrs", "how long PET persists in the environment",
         Inches(9.7), y, num_color=AMBER, border=AMBER)

body_text(sl,
    "PET plastic — bottles, food packaging, polyester fabric — does not break down. "
    "Mechanical recycling only downgrades the material. "
    "Enzymatic degradation breaks PET back to its original building blocks, "
    "which can be repolymerised into virgin-grade plastic.\n\n"
    "That circular loop only works with an enzyme that can handle industrial conditions.",
    t=Inches(4.15), size=18)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — ISPETASE & THE TEMPERATURE WALL
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Background")
slide_title(sl, "IsPETase can degrade plastic \u2014\nbut falls apart above 40 \u00b0C.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

body_text(sl,
    "In 2016, Yoshida et al. discovered Ideonella sakaiensis, a bacterium that lives on PET. "
    "Its enzyme \u2014 IsPETase \u2014 was the first enzyme found to degrade PET at room temperature.\n\n"
    "The problem: IsPETase unfolds above ~40 \u00b0C. "
    "PET chains only become mobile enough for enzymatic attack above the glass transition temperature (~65\u201370 \u00b0C).\n\n"
    "Every degree of thermal stability gained translates directly into faster degradation rates.",
    l=lx, t=BODY_TOP, w=cw, size=18)

# Temperature gap visualisation (right side)
rx2 = Inches(7.2)
box(sl, "Temperature Gap", rx2, BODY_TOP, Inches(5.8), Inches(0.4),
    size=13, bold=True, color=DIM, align=PP_ALIGN.CENTER)

# enzyme zone bar
rect(sl, rx2, Inches(2.5), Inches(2.1), Inches(3.5),
     fill=RGBColor(0x3a, 0x10, 0x10), border=RED, border_pt=1.5)
box(sl, "IsPETase\nStable Zone\n20 \u2013 40 \u00b0C",
    rx2, Inches(3.4), Inches(2.1), Inches(0.9),
    size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

# gap
rect(sl, Inches(9.35), Inches(2.5), Inches(0.7), Inches(3.5),
     fill=RGBColor(0x14, 0x20, 0x30), border=AMBER, border_pt=1)
box(sl, "\u21d5\nGAP", Inches(9.3), Inches(3.7), Inches(0.8), Inches(0.8),
    size=11, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# PET zone bar
rect(sl, Inches(10.1), Inches(2.5), Inches(2.5), Inches(3.5),
     fill=RGBColor(0x06, 0x30, 0x20), border=TEAL, border_pt=1.5)
box(sl, "PET\nReactive Zone\n65 \u2013 80 \u00b0C",
    Inches(10.1), Inches(3.4), Inches(2.5), Inches(0.9),
    size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

box(sl, "\u2190 Close this gap \u2192",
    Inches(9.0), Inches(6.2), Inches(3.5), Inches(0.4),
    size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — OUR APPROACH
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Our Approach")
slide_title(sl, "We built a sequence-only ML model\nto predict which mutations help.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

bullets(sl, [
    "Train on 17,791 real experimental mutations from two curated databases",
    "Predict \u0394\u0394G (change in folding free energy) per mutation",
    "Rank and combine mutations into multi-mutant candidates",
    "Correct scores for real assay conditions: salt, pH, Ca\u00b2\u207a, temperature",
    "Serve results through a public web app \u2014 no software install needed",
], l=lx, t=BODY_TOP, w=cw, size=18)

# Pipeline flow (right side)
steps = [
    ("Protein Sequence (FASTA)", TEAL),
    ("76-Feature Extraction", BLUE),
    ("6-Model Ensemble", BLUE),
    ("Physics Corrections", AMBER),
    ("Ranked Mutation Candidates", TEAL),
]
py = BODY_TOP
pw = Inches(5.5)
for sname, scol in steps:
    rect(sl, rx, py, pw, Inches(0.65), fill=NAVY2, border=scol, border_pt=1.5)
    box(sl, sname, rx + Inches(0.18), py + Inches(0.14),
        pw - Inches(0.36), Inches(0.37), size=16, bold=True, color=WHITE)
    if sname != steps[-1][0]:
        box(sl, "\u25bc", rx + pw / 2 - Inches(0.2), py + Inches(0.65), Inches(0.4), Inches(0.28),
            size=12, color=TEAL, align=PP_ALIGN.CENTER)
    py += Inches(0.95)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — TRAINING DATA
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Training Data")
slide_title(sl, "Only real wet-lab experiments.\nNo synthetic or computationally generated data.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

body_text(sl,
    "We made a deliberate choice: train only on measured \u0394\u0394G values from "
    "wet-lab experiments. Using computationally predicted values would just teach "
    "the model to copy another model\u2019s errors.",
    l=lx, t=BODY_TOP, w=cw, size=18)

card(sl, "FireProtDB",
     "3,438 mutations\nCurated single-point substitutions with experimental \u0394\u0394G and exact assay conditions",
     lx, Inches(3.45), cw, Inches(1.3), border=TEAL, body_size=17)

card(sl, "ThermoMutDB",
     "10,993 mutations across 249 protein families\nTemperature and pH metadata per record",
     lx, Inches(4.9), cw, Inches(1.2), border=BLUE, body_size=17)

# Right: antisymmetry
box(sl, "Thermodynamic Augmentation", rx, BODY_TOP, Inches(5.8), Inches(0.4),
    size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER)

# boxes
rect(sl, rx, Inches(2.6), Inches(2.1), Inches(0.85), fill=DARK, border=TEAL, border_pt=1.5)
box(sl, "S121  (Wild-type)", rx + Inches(0.1), Inches(2.7), Inches(1.9), Inches(0.6),
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, Inches(10.1), Inches(2.6), Inches(2.1), Inches(0.85), fill=DARK, border=BLUE, border_pt=1.5)
box(sl, "E121  (Mutant)", Inches(10.15), Inches(2.7), Inches(1.9), Inches(0.6),
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(sl, "S \u2192 E    \u0394\u0394G = \u22121.2", rx + Inches(2.15), Inches(2.6),
    Inches(1.9), Inches(0.38), size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
box(sl, "E \u2192 S    \u0394\u0394G = +1.2", rx + Inches(2.15), Inches(3.02),
    Inches(1.9), Inches(0.38), size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

rect(sl, rx, Inches(3.68), Inches(5.2), Inches(0.65), fill=NAVY2, border=AMBER, border_pt=1.5)
box(sl, "\u0394\u0394G(A\u2192B) = \u2212\u0394\u0394G(B\u2192A)   \u2014   thermodynamic identity, not an approximation",
    rx + Inches(0.15), Inches(3.74), Inches(4.9), Inches(0.5),
    size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# before / after
rect(sl, rx, Inches(4.6), Inches(2.3), Inches(1.1), fill=NAVY2, border=MID, border_pt=1)
box(sl, "8,905", rx + Inches(0.1), Inches(4.65), Inches(2.1), Inches(0.58),
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(sl, "original measurements", rx + Inches(0.1), Inches(5.18), Inches(2.1), Inches(0.4),
    size=10, color=DIM, align=PP_ALIGN.CENTER)

box(sl, "\u2192", Inches(9.4), Inches(4.8), Inches(0.7), Inches(0.5),
    size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

rect(sl, Inches(10.15), Inches(4.6), Inches(2.3), Inches(1.1), fill=DARK, border=TEAL, border_pt=1.5)
box(sl, "17,791", Inches(10.15), Inches(4.65), Inches(2.3), Inches(0.58),
    size=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
box(sl, "after augmentation", Inches(10.15), Inches(5.18), Inches(2.3), Inches(0.4),
    size=10, color=DIM, align=PP_ALIGN.CENTER)

box(sl, "49.3% stabilising  \u00b7  50.7% destabilising  \u2014  near-perfect balance",
    rx, Inches(5.9), Inches(5.3), Inches(0.4), size=12, color=DIM,
    italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — FEATURES
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Feature Engineering")
slide_title(sl, "Every mutation \u2192 a 76-number vector.\nComputed from sequence alone. No crystal structure.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

bullets(sl, [
    "Physicochemical deltas (12) \u2014 hydrophobicity, volume, charge, helix/sheet propensity, BLOSUM62",
    "Structural context (9) \u2014 estimated RSA, secondary structure fractions, local sequence window",
    "Assay conditions (2) \u2014 temperature and pH from the experimental record",
    "Extended biochemical (18) \u2014 MW, H-bonds, pKa shift, aliphatic index, disorder, Cys distance",
    "Cross-terms (8) \u2014 physically motivated interaction features (e.g. \u0394H \u00d7 temperature)",
], l=lx, t=BODY_TOP, w=cw, size=18)

# Cross-terms cards (right 2x2)
box(sl, "Why Cross-Terms?", rx, BODY_TOP, Inches(5.8), Inches(0.4),
    size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER)

ct = [
    ("\u0394H \u00d7 Temperature", BLUE,
     "Hydrophobic burial stabilises more at high temperature (entropic effect)"),
    ("\u0394Charge \u00d7 Burial", TEAL,
     "Charged residue buried in the core = huge destabilisation (desolvation penalty)"),
    ("\u0394H \u00d7 Burial (RSA)", AMBER,
     "Core mutations have far more impact than surface mutations (packing depth)"),
    ("Burial \u00d7 \u0394Charge(pH)", RED,
     "pH shifts ionisation state; burial amplifies the destabilisation effect"),
]
cw2 = Inches(2.7)
ch2 = Inches(1.75)
for i, (name, col, desc) in enumerate(ct):
    cx = rx if i % 2 == 0 else rx + cw2 + Inches(0.18)
    cy = Inches(2.55) if i < 2 else Inches(2.55) + ch2 + Inches(0.18)
    rect(sl, cx, cy, cw2, ch2, fill=RGBColor(0x0d, 0x1f, 0x35), border=col, border_pt=2)
    box(sl, name, cx + Inches(0.12), cy + Inches(0.12), cw2 - Inches(0.24), Inches(0.35),
        size=12, bold=True, color=col)
    box(sl, desc, cx + Inches(0.12), cy + Inches(0.52), cw2 - Inches(0.24), ch2 - Inches(0.65),
        size=12, color=LIGHT)

box(sl, "We compute these explicitly so the tree models don\u2019t have to discover them from scratch.",
    rx, Inches(6.3), Inches(5.8), Inches(0.5), size=12, color=DIM,
    italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — THE MODEL
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "The Model")
slide_title(sl, "Six independently tuned models.\nThen a meta-classifier stacks them.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

body_text(sl,
    "No single algorithm consistently wins on tabular biological data, "
    "so we train six models and combine them via stacking \u2014 then add a second layer "
    "that directly optimises classification accuracy.",
    l=lx, t=BODY_TOP, w=cw, size=18)

# Table
headers = ["Model", "10-fold CV Accuracy"]
rows = [
    ("GradientBoosting",   "69.2%"),
    ("XGBoost",            "70.4%"),
    ("LightGBM",           "70.1%"),
    ("CatBoost",           "70.5%"),
    ("HistGBM",            "69.9%"),
    ("MLP (neural net)",   "66.9%"),
    ("\u2192 Wide-Stack Ensemble", "71.5%"),
]
col_ws = [Inches(3.2), Inches(2.4)]
row_h  = Inches(0.42)
ty = Inches(3.3)
# header
x = lx
for h, cw3 in zip(headers, col_ws):
    rect(sl, x, ty, cw3, row_h, fill=DARK, border=MID, border_pt=0.5)
    box(sl, h, x + Inches(0.1), ty + Inches(0.07), cw3 - Inches(0.2), row_h - Inches(0.14),
        size=12, bold=True, color=BLUE)
    x += cw3
for rname, racc in rows:
    ty += row_h
    is_win = "\u2192" in rname
    bg = DARK if is_win else NAVY2
    col = TEAL if is_win else LIGHT
    x = lx
    for txt, cw3 in zip([rname, racc], col_ws):
        rect(sl, x, ty, cw3, row_h, fill=bg, border=MID, border_pt=0.5)
        box(sl, txt, x + Inches(0.1), ty + Inches(0.07), cw3 - Inches(0.2),
            row_h - Inches(0.14), size=13, color=col, bold=is_win)
        x += cw3

box(sl, "Hyperparameters tuned with 50-trial Optuna Bayesian search.\nFinal training uses all 17,791 examples.",
    lx, Inches(6.9), Inches(5.7), Inches(0.5), size=11, color=DIM, italic=True)

# Right: stacking diagram
box(sl, "Wide-Stack Architecture", rx, BODY_TOP, Inches(5.8), Inches(0.4),
    size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER)

rect(sl, rx, Inches(2.6), Inches(2.55), Inches(1.55),
     fill=RGBColor(0x0d, 0x1f, 0x35), border=BLUE, border_pt=2)
box(sl, "6 Regressors\nGBM \u00b7 XGB \u00b7 LGBM\nCatBoost \u00b7 HGB \u00b7 MLP\n\u2192 6 OOF \u0394\u0394G values",
    rx + Inches(0.12), Inches(2.65), Inches(2.3), Inches(1.4),
    size=13, color=LIGHT, align=PP_ALIGN.CENTER)

rect(sl, Inches(10.15), Inches(2.6), Inches(2.55), Inches(1.55),
     fill=RGBColor(0x0d, 0x1f, 0x35), border=TEAL, border_pt=2)
box(sl, "3 Direct Classifiers\nXGBClf \u00b7 LGBMClf\nCatBoostClf\n\u2192 3 OOF probabilities",
    Inches(10.2), Inches(2.65), Inches(2.4), Inches(1.4),
    size=13, color=LIGHT, align=PP_ALIGN.CENTER)

box(sl, "\u21e9  combine 9 features  \u21e9",
    rx, Inches(4.25), Inches(5.8), Inches(0.38),
    size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

rect(sl, Inches(8.35), Inches(4.65), Inches(4.35), Inches(0.75),
     fill=DARK, border=AMBER, border_pt=2)
box(sl, "CatBoost Meta-Classifier  (5-fold CV)",
    Inches(8.45), Inches(4.72), Inches(4.15), Inches(0.55),
    size=15, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

box(sl, "\u21e9  threshold sweep [0.30 \u2013 0.70]  \u21e9",
    rx, Inches(5.45), Inches(5.8), Inches(0.35),
    size=13, color=AMBER, align=PP_ALIGN.CENTER)

rect(sl, Inches(8.6), Inches(5.88), Inches(3.9), Inches(0.65),
     fill=RGBColor(0x06, 0x30, 0x20), border=TEAL, border_pt=2)
box(sl, "71.49% CV Accuracy",
    Inches(8.7), Inches(5.93), Inches(3.7), Inches(0.55),
    size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Results")
slide_title(sl, "What the model actually learned.")
divider_line(sl)

# Left: three key stats
big_stat(sl, "71.5%",
         "binary accuracy\n(stabilising vs. destabilising)",
         L, Inches(2.2), w=Inches(3.8), num_color=TEAL, border=TEAL)
big_stat(sl, "r = 0.636",
         "Pearson correlation\npredicted vs. measured \u0394\u0394G",
         L, Inches(3.85), w=Inches(3.8), num_color=BLUE, border=BLUE)
big_stat(sl, "66.1%",
         "LOPO accuracy\n(completely unseen protein families)",
         L, Inches(5.5), w=Inches(3.8), num_color=AMBER, border=AMBER)

# Right: feature importance bar chart
box(sl, "Top Feature Importances  (ensemble average)",
    Inches(5.4), Inches(2.0), Inches(7.7), Inches(0.4),
    size=13, bold=True, color=TEAL)

feats = [
    ("temperature_C",         763, BLUE),
    ("pH",                    733, BLUE),
    ("dH \u00d7 temp  (cross-term)", 652, TEAL),
    ("rel_position",          528, BLUE),
    ("RSA",                   498, BLUE),
    ("nearest_cys_dist",      442, BLUE),
    ("delta_ionization",      429, TEAL),
    ("dAliphatic \u00d7 temp (cross-term)", 388, TEAL),
    ("burial \u00d7 dCharge(pH)  (cross-term)", 360, TEAL),
    ("coil_prop",             350, BLUE),
]
bar_l = Inches(7.0)
bar_max = Inches(5.8)
fy = Inches(2.5)
for fname, fval, fcol in feats:
    box(sl, fname, Inches(5.4), fy, Inches(1.55), Inches(0.32), size=10, color=LIGHT)
    bw = bar_max * fval / 763
    fb = rect(sl, bar_l, fy + Inches(0.04), bw, Inches(0.25),
              fill=fcol, border=None)
    box(sl, str(fval), bar_l + bw + Inches(0.06), fy, Inches(0.55), Inches(0.32),
        size=9, color=DIM)
    fy += Inches(0.38)

box(sl, "\u25a0  Condition / structural features",
    Inches(5.4), Inches(6.45), Inches(3.5), Inches(0.3), size=11, color=BLUE)
box(sl, "\u25a0  Cross-term features",
    Inches(5.4), Inches(6.75), Inches(3.5), Inches(0.3), size=11, color=TEAL)
box(sl, "Temperature and pH rank #1 and #2 \u2014 confirming that including assay conditions was the right call.\nCross-terms appear 4\u00d7 in the top 10, validating the approach of computing interactions explicitly.",
    Inches(5.4), Inches(7.1), Inches(7.7), Inches(0.35), size=10, color=DIM, italic=True)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — PHYSICS CORRECTIONS
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Physics Corrections")
slide_title(sl, "Temperature and pH are in the training data.\nSalt and calcium are not \u2014 so we add explicit physics.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

card(sl, "Debye-H\u00fcckel Correction  (Ionic Strength)",
     "Salt ions screen charge\u2013charge interactions on the protein surface. "
     "At 500\u202fmM NaCl vs. 50\u202fmM, a surface charge mutation matters less because counterions partially cancel it out. "
     "We apply a screening correction only to surface-exposed charged residues (RSA\u202f>\u202f0.25), capped at \u00b10.15 kcal/mol.",
     lx, BODY_TOP, cw, Inches(2.0), border=BLUE, body_size=16)

card(sl, "Ca\u00b2\u207a Hill Equation  (Cutinase Family)",
     "LCC, TfCut2, and Cut190 all have a conserved Ca\u00b2\u207a binding site. "
     "Mutations to chelating residues (Asp/Glu) alter Ca\u00b2\u207a affinity and therefore stability. "
     "We model this with a one-site Hill equation:\n"
     "f = [Ca\u00b2\u207a] / (K\u1d48 + [Ca\u00b2\u207a]),  K\u1d48 = 0.5\u202fmM,  \u0394\u0394G\u2098\u2090\u2093 = 2.0\u202fkcal/mol",
     lx, Inches(4.3), cw, Inches(2.2), border=AMBER, body_size=16)

# Right: visual bars for ionic screening
box(sl, "Ionic Strength Screening Effect",
    rx, BODY_TOP, Inches(5.8), Inches(0.4), size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER)

labels_ionic = [("50 mM\n(lab buffer)", TEAL, 3.5), ("500 mM\n(industrial)", AMBER, 1.6)]
bx = rx
for lbl, col, bar_h in labels_ionic:
    rect(sl, bx + Inches(0.3), Inches(6.05) - Inches(bar_h),
         Inches(1.9), Inches(bar_h), fill=NAVY2, border=col, border_pt=2)
    box(sl, lbl, bx + Inches(0.3), Inches(2.3), Inches(1.9), Inches(0.55),
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    box(sl, "Charge\neffect", bx + Inches(0.3), Inches(6.05) - Inches(bar_h) + Inches(0.1),
        Inches(1.9), Inches(0.6), size=12, color=col, align=PP_ALIGN.CENTER)
    bx += Inches(2.7)

box(sl, "\u2195 ~0.1 kcal/mol difference\nbetween low-salt and high-salt scoring",
    rx, Inches(5.9), Inches(5.5), Inches(0.7), size=13, color=DIM,
    italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 — THE WEB PLATFORM
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Web Platform")
slide_title(sl, "Paste a sequence.\nGet ranked mutation candidates in ~3 seconds.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

bullets(sl, [
    "Condition sliders \u2014 temperature, pH, ionic strength, Ca\u00b2\u207a concentration",
    "Ranked candidate table \u2014 \u0394\u0394G score, predicted \u0394Tm, activity score, literature support",
    "3D molecular viewer \u2014 IsPETase with mutation sites highlighted and an animated guided tour",
    "Literature cross-reference \u2014 auto-checks against FAST-PETase, LCC-ICCG, published variants",
    "Explainability panel \u2014 feature attribution for every prediction",
], l=lx, t=BODY_TOP, w=cw, size=18)

# Code snippet
cb = rect(sl, lx, Inches(5.65), cw, Inches(1.0),
          fill=RGBColor(0x0a, 0x12, 0x20), border=MID, border_pt=1)
box(sl, 'POST /optimize\n{"sequence": "MSGGSSGLPAL...", "target_temperature": 70.0, "ph": 8.0}\n\u2192 ranked candidates in ~3 seconds',
    lx + Inches(0.15), Inches(5.7), cw - Inches(0.3), Inches(0.9), size=11, color=TEAL)

# Right: layer diagram
box(sl, "System Architecture", rx, BODY_TOP, Inches(5.8), Inches(0.4),
    size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER)

arch = [
    ("React Front-End", BLUE,
     "Sequence input  \u00b7  Condition sliders  \u00b7  3Dmol.js viewer"),
    ("FastAPI Backend", TEAL,
     "/optimize  \u00b7  /explain  \u00b7  /literature  \u00b7  /classifier"),
]
py = Inches(2.55)
for aname, acol, asub in arch:
    rect(sl, rx, py, Inches(5.8), Inches(0.85), fill=NAVY2, border=acol, border_pt=2)
    box(sl, aname, rx + Inches(0.2), py + Inches(0.07),
        Inches(5.4), Inches(0.35), size=15, bold=True, color=WHITE)
    box(sl, asub, rx + Inches(0.2), py + Inches(0.46),
        Inches(5.4), Inches(0.32), size=11, color=DIM)
    py += Inches(1.05)

# Three model boxes
names_mc = [("Ensemble\n6 regressors + meta-clf", BLUE),
            ("\u0394Tm Model\nMAE = 4.6 \u00b0C", TEAL),
            ("Literature DB\nFAST-PETase, LCC-ICCG", AMBER)]
bx = rx
for nm, nc in names_mc:
    rect(sl, bx, Inches(4.75), Inches(1.85), Inches(1.05),
         fill=NAVY2, border=nc, border_pt=1.5)
    box(sl, nm, bx + Inches(0.1), Inches(4.8), Inches(1.65), Inches(0.95),
        size=12, color=LIGHT, align=PP_ALIGN.CENTER)
    bx += Inches(1.97)

rect(sl, rx, Inches(5.95), Inches(5.8), Inches(0.75),
     fill=NAVY2, border=MID, border_pt=1)
box(sl, "3Dmol.js  \u2014  IsPETase (PDB: 5XJH)  \u00b7  Geometry baked once \u2192 pure camera moves = 60\u202ffps",
    rx + Inches(0.2), Inches(6.0), Inches(5.4), Inches(0.65),
    size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 — LIMITATIONS
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Honest Assessment")
slide_title(sl, "What 71.5% accuracy means \u2014\nand what could get us to 80%.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

box(sl, "Current limitations:", lx, BODY_TOP, cw, Inches(0.4),
    size=16, bold=True, color=WHITE)
bullets(sl, [
    "No structural features \u2014 FoldX/Rosetta use explicit 3D contacts; our RSA is sequence-estimated",
    "No PSSM profiles \u2014 evolutionary conservation adds ~2\u20134% but takes minutes per protein",
    "Additivity assumption \u2014 \u0394\u0394G values are summed for multi-mutants, which breaks down for nearby residues",
    "LOPO gap \u2014 66.1% on unseen protein families vs. 71.5% overall",
], l=lx, t=Inches(2.55), w=cw, size=17)

box(sl, "Path to 80%:", rx, BODY_TOP, Inches(5.8), Inches(0.4),
    size=16, bold=True, color=WHITE)

card(sl, "AlphaFold2 Predicted RSA",
     "Replace sequence-estimated RSA with AlphaFold2 predicted structures.\nEstimated gain: +3\u20134% accuracy.",
     rx, Inches(2.55), Inches(5.8), Inches(1.2), border=TEAL, body_size=16)

card(sl, "PSSM Conservation Features",
     "Pre-computed BLAST profiles for common enzyme families.\nEstimated gain: +2\u20133%. Already partially implemented.",
     rx, Inches(3.9), Inches(5.8), Inches(1.2), border=BLUE, body_size=16)

card(sl, "ProDDG / S2648 Dataset",
     "The S2648 dataset was not used in our current model. "
     "Adding it would provide ~2,600 more training examples across new protein families.\nEstimated gain: +1\u20132%.",
     rx, Inches(5.25), Inches(5.8), Inches(1.2), border=AMBER, body_size=16)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Conclusion")
slide_title(sl, "Sequence in.  Ranked candidates out.\nNo crystal structure. No software to install.")
divider_line(sl)

lx, rx, cw = two_col_split(sl)

body_text(sl,
    "We built a complete ML pipeline for protein stability prediction "
    "\u2014 data curation, physics-informed features, stacking ensemble, physics corrections, and a public web app "
    "\u2014 all from sequence alone.\n\n"
    "71.5% binary accuracy and r\u202f=\u202f0.636 Pearson correlation are competitive with published sequence-only "
    "\u0394\u0394G predictors. More importantly, the system is practical: "
    "cut 10,000 possible single-point mutations down to 50 high-confidence candidates for wet-lab validation.",
    l=lx, t=BODY_TOP, w=cw, size=18)

big_stat(sl, "17,791", "real experimental mutations in training set",
         rx, Inches(2.1), w=Inches(5.5), num_color=WHITE, border=MID)
big_stat(sl, "71.5%", "binary classification accuracy (stacking CV)",
         rx, Inches(3.6), w=Inches(5.5), num_color=TEAL, border=TEAL)
big_stat(sl, "~3 sec", "to rank all candidates for a new protein sequence",
         rx, Inches(5.1), w=Inches(5.5), num_color=BLUE, border=BLUE)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 — ACKNOWLEDGMENTS
# ═══════════════════════════════════════════════════════════
sl = new_slide()
slide_label(sl, "Acknowledgments & References")
slide_title(sl, "Thank you.", size=52)
divider_line(sl, top=Inches(1.75))

lx, rx, cw = two_col_split(sl)

box(sl, "Databases that made this work possible:",
    lx, Inches(2.05), cw, Inches(0.4), size=16, bold=True, color=WHITE)
bullets(sl, [
    "FireProtDB \u2014 Stourac et al., curated experimental thermostability data",
    "ThermoMutDB \u2014 Pucci et al. (2021), Nucleic Acids Research",
    "RCSB Protein Data Bank \u2014 Berman et al. (2000)",
    "3Dmol.js \u2014 3D structure visualisation library",
], l=lx, t=Inches(2.5), w=cw, size=17)

box(sl, "Key prior work:",
    lx, Inches(4.5), cw, Inches(0.4), size=16, bold=True, color=WHITE)
bullets(sl, [
    "Lu et al. (2022) \u2014 FAST-PETase, Nature",
    "Tournier et al. (2020) \u2014 LCC-ICCG, Nature",
    "Yoshida et al. (2016) \u2014 IsPETase discovery, Science",
], l=lx, t=Inches(4.95), w=cw, size=17)

box(sl, "Questions?",
    rx, Inches(3.0), Inches(5.8), Inches(1.0),
    size=42, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

box(sl, "PETase-ML  \u00b7  Los Altos Hacks 2026",
    rx, Inches(4.2), Inches(5.8), Inches(0.45),
    size=16, color=DIM, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
out = '/Users/admin/Documents/Los-Altos-Hacks/PETase_ML_Presentation.pptx'
prs.save(out)
print(f'Done: {out}')
